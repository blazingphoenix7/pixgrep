"""Tests for POST /api/export/pptx endpoint."""
from __future__ import annotations

import io

import pytest
from fastapi.testclient import TestClient
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from pixgrep.server import create_app


# ── Minimal fake engine ───────────────────────────────────────────────────────

class _FakeTags:
    has_data = False
    def facets(self): return {}


class _FakeEngine:
    def __init__(self, index_dir, paths):
        self.index_dir = index_dir
        self._paths = list(paths)
        self.count = len(self._paths)
        self._tags = _FakeTags()

    def path_for(self, row: int) -> str:
        if row < 0 or row >= len(self._paths):
            raise IndexError(row)
        return str(self._paths[row])


# ── Helpers ───────────────────────────────────────────────────────────────────

def _png(path, size=(60, 40), color=(180, 120, 60)):
    Image.new("RGB", size, color=color).save(path, format="PNG")


def _tiff(path, size=(60, 40), color=(60, 120, 180)):
    Image.new("RGB", size, color=color).save(path, format="TIFF")


def _pics(slide):
    return [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]


def _textboxes(slide):
    return [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.TEXT_BOX]


# ── Fixtures ──────────────────────────────────────────────────────────────────

@pytest.fixture
def client(tmp_path):
    imgs = []
    for i in range(5):
        p = tmp_path / f"img_{i}.png"
        _png(p)
        imgs.append(p)
    return TestClient(create_app(_FakeEngine(tmp_path, imgs)))


@pytest.fixture
def tiff_client(tmp_path):
    p = tmp_path / "sample.tiff"
    _tiff(p)
    return TestClient(create_app(_FakeEngine(tmp_path, [p])))


# ── Tests ─────────────────────────────────────────────────────────────────────

def test_valid_pptx_layout1(client):
    r = client.post("/api/export/pptx", json={"rows": [0, 1, 2], "layout": "1", "captions": True})
    assert r.status_code == 200
    assert "presentationml" in r.headers["content-type"]
    prs = Presentation(io.BytesIO(r.content))
    assert len(prs.slides) == 3
    for slide in prs.slides:
        assert len(_pics(slide)) == 1


def test_layout4_five_images_two_slides(client):
    r = client.post("/api/export/pptx", json={"rows": [0, 1, 2, 3, 4], "layout": "4", "captions": False})
    assert r.status_code == 200
    prs = Presentation(io.BytesIO(r.content))
    assert len(prs.slides) == 2
    # first slide has 4 pictures, second has 1
    assert len(_pics(prs.slides[0])) == 4
    assert len(_pics(prs.slides[1])) == 1


def test_captions_on_adds_textboxes(client):
    r = client.post("/api/export/pptx", json={"rows": [0], "layout": "1", "captions": True})
    assert r.status_code == 200
    prs = Presentation(io.BytesIO(r.content))
    assert len(_textboxes(prs.slides[0])) == 1


def test_captions_off_no_textboxes(client):
    r = client.post("/api/export/pptx", json={"rows": [0], "layout": "1", "captions": False})
    assert r.status_code == 200
    prs = Presentation(io.BytesIO(r.content))
    assert len(_textboxes(prs.slides[0])) == 0


def test_empty_rows_returns_422(client):
    r = client.post("/api/export/pptx", json={"rows": [], "layout": "1", "captions": True})
    assert r.status_code == 422


def test_too_many_rows_returns_413(client):
    r = client.post("/api/export/pptx", json={"rows": list(range(201)), "layout": "1", "captions": False})
    assert r.status_code == 413


def test_unknown_row_skipped(client):
    r = client.post("/api/export/pptx", json={"rows": [0, 9999], "layout": "1", "captions": False})
    assert r.status_code == 200
    prs = Presentation(io.BytesIO(r.content))
    assert len(prs.slides) == 1


def test_all_unknown_rows_returns_422(client):
    r = client.post("/api/export/pptx", json={"rows": [9998, 9999], "layout": "1", "captions": False})
    assert r.status_code == 422


def test_tiff_embeds_via_conversion(tiff_client):
    r = tiff_client.post("/api/export/pptx", json={"rows": [0], "layout": "1", "captions": False})
    assert r.status_code == 200
    prs = Presentation(io.BytesIO(r.content))
    assert len(prs.slides) == 1
    assert len(_pics(prs.slides[0])) == 1
