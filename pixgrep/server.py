from __future__ import annotations

import io
from pathlib import Path
from typing import Literal, Optional

from fastapi import FastAPI, File, HTTPException, Query, UploadFile
from fastapi.responses import FileResponse, Response
from fastapi.staticfiles import StaticFiles
from PIL import Image, UnidentifiedImageError
from starlette.middleware.base import BaseHTTPMiddleware

from pydantic import BaseModel

from .search import SearchEngine


class _CacheControlMiddleware(BaseHTTPMiddleware):
    async def dispatch(self, request, call_next):
        response = await call_next(request)
        path = request.url.path
        if path.startswith("/api/"):
            response.headers["Cache-Control"] = "no-cache"
        elif path == "/" or path == "/index.html" or path.startswith("/static/"):
            response.headers["Cache-Control"] = "no-store"
        return response

STATIC_DIR = Path(__file__).resolve().parent.parent / "static"


def _parse_filters(f_params: list[str]) -> dict[str, str] | None:
    """Parse repeated `f=field:value` query params into a filters dict."""
    if not f_params:
        return None
    filters: dict[str, str] = {}
    for item in f_params:
        if ":" in item:
            field, _, value = item.partition(":")
            field = field.strip()
            value = value.strip()
            if field and value:
                filters[field] = value
    return filters if filters else None


class _ExportRequest(BaseModel):
    rows: list[int]
    layout: Literal["1", "4"] = "1"
    captions: bool = True


def _prepare_image(path: Path) -> tuple:
    """Returns (pptx_src, width_px, height_px, stem). Converts exotic formats to JPEG in-memory."""
    img = Image.open(path)
    w, h = img.size
    if path.suffix.lower() in (".jpg", ".jpeg", ".png"):
        return str(path), w, h, path.stem
    buf = io.BytesIO()
    img.convert("RGB").save(buf, format="JPEG")
    buf.seek(0)
    return buf, w, h, path.stem


def _build_pptx(prepared: list, layout: str, captions: bool) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    SW, SH = Inches(13.333), Inches(7.5)
    M, CAP_H, GAP = Inches(0.4), Inches(0.25), Inches(0.2)

    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    blank = prs.slide_layouts[6]

    def _fit(iw: int, ih: int, cw: int, ch: int) -> tuple:
        if not (iw and ih):
            return cw, ch
        ar = iw / ih
        if ar * ch >= cw:
            return cw, max(1, int(cw / ar))
        return max(1, int(ch * ar)), ch

    def _place(slide, src, iw: int, ih: int, x: int, y: int, cw: int, ch: int, stem: str) -> None:
        img_h = ch - (CAP_H if captions else 0)
        fw, fh = _fit(iw, ih, cw, img_h)
        left = x + (cw - fw) // 2
        top = y + (img_h - fh) // 2
        if isinstance(src, io.BytesIO):
            src.seek(0)
        slide.shapes.add_picture(src, left, top, fw, fh)
        if captions:
            tb = slide.shapes.add_textbox(x, y + img_h, cw, CAP_H)
            tf = tb.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = stem
            run.font.size = Pt(10)
            run.font.color.rgb = RGBColor(0x8B, 0x93, 0xA1)

    if layout == "1":
        for src, iw, ih, stem in prepared:
            slide = prs.slides.add_slide(blank)
            _place(slide, src, iw, ih, M, M, SW - 2 * M, SH - 2 * M, stem)
    else:
        cw = (SW - 2 * M - GAP) // 2
        ch = (SH - 2 * M - GAP) // 2
        for i in range(0, len(prepared), 4):
            slide = prs.slides.add_slide(blank)
            for j, (src, iw, ih, stem) in enumerate(prepared[i:i + 4]):
                col, row_i = j % 2, j // 2
                x = M + col * (cw + GAP)
                y = M + row_i * (ch + GAP)
                _place(slide, src, iw, ih, x, y, cw, ch, stem)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def create_app(engine: SearchEngine) -> FastAPI:
    app = FastAPI(title="pixgrep")
    app.add_middleware(_CacheControlMiddleware)

    @app.get("/api/meta")
    def meta():
        return {"count": engine.count}

    @app.get("/api/facets")
    def facets():
        if not engine._tags.has_data:
            return {}
        raw = engine._tags.facets()
        return {
            field: [
                {"value": v, "count": c}
                for v, c in sorted(val_counts.items(), key=lambda x: -x[1])
            ]
            for field, val_counts in raw.items()
        }

    @app.get("/api/search")
    def search(
        q: str = Query(...),
        k: int = Query(24, ge=1, le=200),
        min_ratio: float = Query(0.6, ge=0.0, le=1.0),
        min_score: float = Query(0.05, ge=0.0, le=1.0),
        f: list[str] = Query(default=[]),
        hw: Optional[float] = Query(default=None, ge=0.0, le=1.0),
    ):
        if not q.strip():
            raise HTTPException(status_code=400, detail="empty query")
        filters = _parse_filters(f)
        results = engine.text_search(
            q, k=k, min_ratio=min_ratio, min_score=min_score,
            filters=filters, hybrid_weight=hw,
        )
        return {"results": results}

    @app.post("/api/search/image")
    def search_image(
        file: UploadFile = File(...),
        k: int = Query(24, ge=1, le=200),
        min_ratio: float = Query(0.6, ge=0.0, le=1.0),
        min_score: float = Query(0.05, ge=0.0, le=1.0),
        f: list[str] = Query(default=[]),
    ):
        data = file.file.read()
        try:
            img = Image.open(io.BytesIO(data)).convert("RGB")
        except (UnidentifiedImageError, OSError):
            raise HTTPException(status_code=400, detail="could not decode image")
        filters = _parse_filters(f)
        results = engine.image_search(
            img, k=k, min_ratio=min_ratio, min_score=min_score, filters=filters,
        )
        return {"results": results}

    @app.get("/api/similar/{row}")
    def similar(
        row: int,
        k: int = Query(24, ge=1, le=200),
        min_ratio: float = Query(0.6, ge=0.0, le=1.0),
        min_score: float = Query(0.05, ge=0.0, le=1.0),
        f: list[str] = Query(default=[]),
    ):
        filters = _parse_filters(f)
        try:
            results = engine.similar(
                row, k=k, min_ratio=min_ratio, min_score=min_score, filters=filters,
            )
        except IndexError:
            raise HTTPException(status_code=404, detail="unknown row")
        return {"results": results}

    @app.get("/api/group/{row}")
    def group(row: int):
        try:
            results = engine.group_members(row)
        except IndexError:
            raise HTTPException(status_code=404, detail="unknown row")
        return {"results": results}

    @app.get("/api/image/{row}")
    def image(row: int):
        try:
            path = Path(engine.path_for(row))
        except IndexError:
            raise HTTPException(status_code=404, detail="unknown row")
        if not path.is_file():
            raise HTTPException(status_code=404, detail="file missing on disk")
        return FileResponse(path)

    @app.get("/api/thumb/{row}")
    def thumb(row: int):
        thumb_path = engine.index_dir / "thumbs" / f"{row}.jpg"
        if thumb_path.is_file():
            return FileResponse(thumb_path)
        try:
            path = Path(engine.path_for(row))
        except IndexError:
            raise HTTPException(status_code=404, detail="unknown row")
        if not path.is_file():
            raise HTTPException(status_code=404, detail="file missing on disk")
        return FileResponse(path)

    @app.post("/api/export/pptx")
    def export_pptx(req: _ExportRequest):
        if len(req.rows) > 200:
            raise HTTPException(status_code=413, detail="too many rows (max 200)")
        if not req.rows:
            raise HTTPException(status_code=422, detail="rows list is empty")
        prepared = []
        for row in req.rows:
            try:
                path = Path(engine.path_for(row))
            except IndexError:
                continue
            if not path.is_file():
                continue
            try:
                prepared.append(_prepare_image(path))
            except Exception:
                continue
        if not prepared:
            raise HTTPException(status_code=422, detail="no valid images after filtering")
        pptx_bytes = _build_pptx(prepared, req.layout, req.captions)
        return Response(
            content=pptx_bytes,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": 'attachment; filename="pixgrep-export.pptx"'},
        )

    if STATIC_DIR.is_dir():
        app.mount("/static", StaticFiles(directory=STATIC_DIR), name="static")

        @app.get("/")
        def root():
            return FileResponse(STATIC_DIR / "index.html")

    return app


def main() -> None:
    import argparse

    import uvicorn

    from .config import load_config

    parser = argparse.ArgumentParser(description="pixgrep search server")
    parser.add_argument("--host", default="127.0.0.1")
    parser.add_argument("--port", type=int, default=8177)
    args = parser.parse_args()

    cfg = load_config()
    print(f"Loading model {cfg.model_id} ...")
    engine = SearchEngine(
        cfg.index_dir,
        cfg.make_embedder(),
        hybrid_weight=cfg.hybrid_weight,
        junk_threshold=cfg.junk_threshold,
        group_strip_pattern=cfg.group_strip_pattern,
        near_dupe_cos=cfg.near_dupe_cos,
        lexical_inject_k=cfg.lexical_inject_k,
        junk_soft_weight=cfg.junk_soft_weight,
    )
    print(f"Index loaded: {engine.count} images. http://{args.host}:{args.port}")
    uvicorn.run(create_app(engine), host=args.host, port=args.port)


if __name__ == "__main__":
    main()
